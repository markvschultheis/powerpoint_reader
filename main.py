from pptx import Presentation
from tkinter import *
from tkinter.ttk import *

from tkinter.filedialog import askopenfile, asksaveasfile

root = Tk()
root.geometry('250x100')
root.title('Church Setup')


def open_file():
    try:
        file = askopenfile()
        file_string = str(file)
        file_string = file_string[file_string.find('name=') + 6: file_string.find('mode=') - 2]
        read_powerpoint(file_string)
        btn.configure(text="Save File")
        btn.configure(command=save_file)
    except:
        print('hi')
        btn.configure(text="Open Powerpoint File")


def read_powerpoint(powerpoint_path):
    prs = Presentation(powerpoint_path)
    text_file = open("text.txt", 'w')
    try:
        values = open('values.txt', 'r')
        hymn_values = []
        for index in range(3):
            hymn_values.append(str_to_bool(values.readline()))
        fr_bool = str_to_bool(values.readline())
        sr_bool = str_to_bool(values.readline())
        a_bool = str_to_bool(values.readline())
    except:
        values = open('values.txt', 'w')
        values.close()
        values = open('values.txt', 'r')
        hymn_values = []
        fr_bool = False
        sr_bool = False
        a_bool = False

    text_runs = []
    text_string = ''
    count = 0
    hymn_number = 1
    gospel_count = 1
    need_number_verses = False
    first_reading_going = False
    second_reading_going = False
    possible_hymn_start = False
    hymns = []
    hymn_counts = []
    hymn_multiple = []
    hymn_next_slide = 0
    r1 = -1
    r1e = -1
    r2 = -1
    r2e = -1
    a = -1
    g1 = -1
    g2 = -1
    first_reading_possible = False
    first_reading_done = False
    second_reading_possible = False
    second_reading_done = False
    possible_number = False
    numbers = {
        1: 'first',
        2: 'second',
        3: 'third',
        4: 'fourth'
    }
    service_end = -1
    text_string += 'dim streaming As Boolean = True\n\n'
    for slide in prs.slides:
        count += 1
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    print(run.text)
                    text_runs.append(run.text)
                    if need_number_verses:
                        hymn_counts.append(run.text)
                        need_number_verses = False
                    elif first_reading_possible:
                        if 'reading' in run.text.lower() or 'testament' in run.text.lower():
                            r1 = count
                            first_reading_going = True
                        first_reading_possible = False
                    elif second_reading_possible:
                        if 'reading' in run.text.lower():
                            r2 = count
                            second_reading_going = True
                        second_reading_possible = False
                    elif possible_hymn_start:
                        if '/' in run.text:
                            hymns.append(count)
                            hymn_number += 1
                            value = run.text[run.text.find('/') + 1:]
                            if value != '':
                                hymn_counts.append(run.text[run.text.find('/') + 1:])
                                hymn_next_slide = count + 1
                            else:
                                possible_number = True
                        possible_hymn_start = False
                    elif possible_number:
                        if run.text in ['1', '2', '3', '4', '5', '6', '7', '8', '9']:
                            hymn_counts.append(run.text)
                            hymn_next_slide = count + 1
                        possible_number = False
                    elif first_reading_going:
                        if 'this is the word of the lord' in run.text.lower():
                            r1e = count
                            first_reading_going = False
                            first_reading_done = True
                    elif second_reading_going:
                        if 'this is the word of the lord' in run.text.lower():
                            r2e = count
                            second_reading_going = False
                            second_reading_done = True
                    elif run.text == '1/' and count != hymn_next_slide:
                        hymns.append(count)
                        hymn_number += 1
                        need_number_verses = True
                        hymn_next_slide = count + 1

                    elif (
                            (
                                    'first reading' in run.text.lower() or 'old testament' in run.text.lower()) and not first_reading_done) and fr_bool:
                        r1 = count
                        first_reading_going = True
                    elif ((
                                  'first' in run.text.lower() or 'old' in run.text.lower()) and not first_reading_done) and fr_bool:
                        first_reading_possible = True
                    elif ((
                                  'epistle' in run.text.lower() or 'second reading' in run.text.lower()) and not second_reading_done) and sr_bool:
                        r2 = count
                        second_reading_going = True
                    elif ('second' in run.text.lower() and not second_reading_done) and sr_bool:
                        second_reading_possible = True
                    elif '1' in run.text:
                        if count == hymn_next_slide:

                            hymn_multiple.append(2)
                            hymn_next_slide = count + 1

                        elif not first_reading_going or not second_reading_going:
                            possible_hymn_start = True
                            hymn_next_slide = 0
                    elif count == hymn_next_slide:
                        hymn_multiple.append(1)
                        hymn_next_slide = 0
                    elif 'anthem' in run.text.lower() and a_bool:
                        a = count
                    elif 'gospel verse' in run.text.lower():
                        if gospel_count == 1:
                            g1 = count
                        else:
                            g2 = count
                        gospel_count += 1
                    elif 'preface' in run.text.lower():
                        service_end = count

    for var in range(3):
        if not hymn_values[var]:
            hymns[var] = -1
            hymn_counts[var] = -1

    for var in range(4 - len(hymns)):
        hymns.append(-1)
        hymn_counts.append(-1)
    for var in range(4 - len(hymn_multiple)):
        hymn_multiple.append(1)
    count = 1
    for start, number, multiple in zip(hymns, hymn_counts, hymn_multiple):
        text_string += 'dim {}HymnStart As Integer = {}\n'.format(numbers[count], str(start))

        text_string += 'dim {}HymnNumberOfVerses As Integer = {}\n\n'.format(numbers[count],
                                                                             str(int(number) * int(multiple)))
        count += 1

    text_string += 'dim firstReadingStart As Integer = ' + str(r1) + '\n'
    text_string += "dim firstReadingEnd As Integer = " + str(r1e) + '\n\n'
    text_string += "dim anthemStart As Integer = " + str(a) + '\n\n'
    text_string += 'dim secondReadingStart As Integer = ' + str(r2) + '\n'
    text_string += "dim secondReadingEnd As Integer = " + str(r2e) + '\n\n'
    text_string += "dim " + numbers[1] + 'GospelVerse As Integer = ' + str(g1) + '\n'
    text_string += "dim " + numbers[2] + 'GospelVerse As Integer = ' + str(g2) + '\n\n'
    text_string += "dim serviceEnd As Integer = " + str(service_end) + '\n\n'
    text_file.write(text_string)
    text_file.close()
    values.close()


def save_file():
    text_file = open("text.txt", 'r')
    end = open('end_of_script.txt', 'r')
    try:
        file = asksaveasfile(filetypes=(("Text Document", "*.txt"),
                                        ('All Files', '*.*')),
                             title='Save File',
                             initialfile="Script1.txt",

                             )
        text_string = ''

        lines = text_file.readlines()
        for line in lines:
            text_string += str(line)
        lines = end.readlines()
        for line in lines:
            text_string += str(line)
        file.write(text_string)
        file.close()
        text_file.close()
        end.close()
        btn.configure(text="Open Powerpoint File")
        btn.configure(command=open_file)
    except:
        btn.configure(text="Save File")


def check_boxes():
    hymn1_check.pack(anchor=NW, padx=70)
    hymn2_check.pack(anchor=NW, padx=70)
    hymn3_check.pack(anchor=NW, padx=70)
    fr_check.pack(anchor=NW, padx=70)
    sr_check.pack(anchor=NW, padx=70)
    a_check.pack(anchor=NW, padx=70)
    save_values.pack(pady=10)
    root.geometry('250x250')


def get_values():
    hymn1_check.pack_forget()
    hymn2_check.pack_forget()
    hymn3_check.pack_forget()
    fr_check.pack_forget()
    sr_check.pack_forget()
    a_check.pack_forget()
    save_values.pack_forget()
    root.geometry('250x100')
    values = open('values.txt', 'w')
    values.write(bool_to_str(h1c.get()))
    values.write(bool_to_str(h2c.get()))
    values.write(bool_to_str(h3c.get()))
    values.write(bool_to_str(fr.get()))
    values.write(bool_to_str(sr.get()))
    values.write(bool_to_str(a.get()))
    values.close()


def bool_to_str(value):
    if value:
        return 'True\n'
    else:
        return 'False\n'


def str_to_bool(value):
    if 'true' in value.lower():
        return True
    else:
        return False


h1c = BooleanVar()
hymn1_check = Checkbutton(root, text='First Hymn', variable=h1c, onvalue=True, offvalue=False)
h2c = BooleanVar()
hymn2_check = Checkbutton(root, text='Second Hymn', variable=h2c, onvalue=True, offvalue=False)
h3c = BooleanVar()
hymn3_check = Checkbutton(root, text='Third Hymn', variable=h3c, onvalue=True, offvalue=False)
fr = BooleanVar()
fr_check = Checkbutton(root, text='First Reading', variable=fr, onvalue=True, offvalue=False)
sr = BooleanVar()
sr_check = Checkbutton(root, text='Second Reading', variable=sr, onvalue=True, offvalue=False)
a = BooleanVar()
a_check = Checkbutton(root, text='Anthem', variable=a, onvalue=True, offvalue=False)
save_values = Button(root, text='Save Values', command=get_values)
btn = Button(root, text='Open Powerpoint File', command=open_file)
btn.pack(side=TOP, pady=10)
check = Button(root, text='Check Uploads', command=check_boxes).pack()

mainloop()
